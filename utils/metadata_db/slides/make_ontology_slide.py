#!/usr/bin/env python3
"""
Generate a single PowerPoint slide summarising the "ontology-guided retrieval"
idea for AngioVision DSA image similarity.

    python make_ontology_slide.py            # -> ontology_retrieval_slide.pptx

Pure python-pptx; the arterial tree and all graphics are drawn as native vector
shapes (fully editable in PowerPoint/Keynote).  Accent colours match the existing
AngioVision figures (eval_kn_retrieval.py / eda_labeled_sequences.py).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ── Palette (reused from AngioVision figures) ────────────────────────────────
NAVY   = RGBColor(0x2C, 0x3E, 0x50)
BLUE   = RGBColor(0x29, 0x80, 0xB9)
TEAL   = RGBColor(0x5D, 0xAD, 0xE2)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
AMBER  = RGBColor(0xF3, 0x9C, 0x12)
AMBERD = RGBColor(0xB9, 0x77, 0x0E)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
GREY   = RGBColor(0x57, 0x63, 0x6D)
LIGHT  = RGBColor(0xEC, 0xF1, 0xF6)
CARD   = RGBColor(0xF7, 0xF9, 0xFC)
BORDER = RGBColor(0xC8, 0xD2, 0xDC)
SUBTLE = RGBColor(0xC8, 0xD6, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2C, 0x3E, 0x50)

FONT = "Calibri"


# ── Helpers ──────────────────────────────────────────────────────────────────
def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0,
        shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(shape_or_slide, runs, l=None, t=None, w=None, h=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, space_after=2,
         line_spacing=1.0, wrap=True, mleft=0.06):
    """`runs` = list of paragraphs; each paragraph = list of (str, fmt-dict)."""
    if l is not None:  # create a textbox
        shape = shape_or_slide.shapes.add_textbox(
            Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        shape = shape_or_slide
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(mleft)
    tf.margin_right = Inches(mleft)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for txt, fmt in para_runs:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = fmt.get("name", FONT)
            f.size = Pt(fmt.get("size", 12))
            f.bold = fmt.get("bold", False)
            f.italic = fmt.get("italic", False)
            if fmt.get("color") is not None:
                f.color.rgb = fmt["color"]
    return shape


def node(slide, l, t, label, fill, txt=WHITE, w=1.05, h=0.42, fs=10.5):
    sp = box(slide, l, t, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.28)
    text(sp, [[(label, {"size": fs, "bold": True, "color": txt})]],
         align=PP_ALIGN.CENTER, space_after=0)
    return sp


def connector(slide, x1, y1, x2, y2, color=GREY, w=1.5):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    _no_shadow(c)
    return c


def card(slide, l, t, w, h, letter, accent, lead, body):
    box(slide, l, t, w, h, fill=CARD, line=BORDER, line_w=0.75,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    d = 0.44
    bd = box(slide, l + 0.16, t + (h - d) / 2, d, d, fill=accent,
             shape=MSO_SHAPE.OVAL)
    text(bd, [[(letter, {"size": 15, "bold": True, "color": WHITE})]],
         align=PP_ALIGN.CENTER, space_after=0)
    text(slide,
         [[(lead + "  ", {"size": 12.5, "bold": True, "color": DARK}),
           (body, {"size": 11.5, "color": GREY})]],
         l=l + 0.74, t=t, w=w - 0.86, h=h, align=PP_ALIGN.LEFT,
         line_spacing=1.0, space_after=0)


# ── Build slide ───────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    box(slide, 0, 0, 13.333, 1.02, fill=NAVY)
    box(slide, 0, 1.02, 13.333, 0.06, fill=BLUE)
    text(slide, [[("Ontology-Guided Retrieval for Deep Subtraction Angiography",
                   {"size": 22, "bold": True, "color": WHITE})]],
         l=0.45, t=0.08, w=12.4, h=0.58, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    text(slide, [[("Injecting an artery ontology into the RAD-DINO + ChromaDB "
                   "image-similarity pipeline",
                   {"size": 12.5, "italic": True, "color": SUBTLE})]],
         l=0.47, t=0.62, w=12.4, h=0.34, anchor=MSO_ANCHOR.TOP, wrap=False)

    # Column headers + divider
    text(slide, [[("The Problem", {"size": 16, "bold": True, "color": RED})]],
         l=0.45, t=1.18, w=5.7, h=0.4)
    text(slide, [[("The Idea \u2014 inject an artery ontology",
                   {"size": 16, "bold": True, "color": BLUE})]],
         l=6.55, t=1.18, w=6.5, h=0.4)
    connector(slide, 6.35, 1.5, 6.35, 6.08, color=BORDER, w=1.0)

    # ── LEFT: problem framing + arterial tree ────────────────────────────────
    text(slide, [[("RAD-DINO + ChromaDB kNN classifies 24 fine-grained artery "
                   "labels, scored by exact match.",
                   {"size": 11.5, "color": GREY})]],
         l=0.45, t=1.56, w=5.75, h=0.5, anchor=MSO_ANCHOR.TOP)

    # subtree highlight (drawn first so it sits behind nodes)
    box(slide, 0.55, 2.66, 4.45, 2.84, line=AMBER, line_w=2.25,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)

    # connectors
    connector(slide, 3.075, 2.57, 3.075, 2.78)   # Aorta -> Celiac
    connector(slide, 3.075, 3.20, 1.875, 3.50)   # Celiac -> CHA
    connector(slide, 3.075, 3.20, 4.375, 3.50)   # Celiac -> Splenic
    connector(slide, 1.875, 3.92, 1.875, 4.22)   # CHA -> PHA
    connector(slide, 1.875, 4.64, 1.225, 4.94)   # PHA -> HA-L
    connector(slide, 1.875, 4.64, 2.525, 4.94)   # PHA -> HA-R

    # nodes
    node(slide, 2.55, 2.15, "Aorta", NAVY)
    node(slide, 2.55, 2.78, "Celiac trunk", BLUE)
    node(slide, 1.35, 3.50, "Common hepatic", TEAL, txt=DARK)
    node(slide, 3.85, 3.50, "Splenic", TEAL, txt=DARK)
    node(slide, 1.35, 4.22, "Proper hepatic", TEAL, txt=DARK)
    node(slide, 0.70, 4.94, "Hepatic L", TEAL, txt=DARK)
    node(slide, 2.00, 4.94, "Hepatic R", TEAL, txt=DARK)

    text(slide, [[("One selective injection \u2192 the whole subtree opacifies",
                   {"size": 10.5, "bold": True, "color": AMBERD})]],
         l=0.55, t=5.5, w=4.45, h=0.3, align=PP_ALIGN.CENTER)

    text(slide,
         [[("Exact-match penalises anatomically-correct hits: ",
            {"size": 10.5, "color": RED}),
           ("CT\u2192CHA", {"size": 10.5, "bold": True, "color": RED}),
           (" (parent/child), ", {"size": 10.5, "color": RED}),
           ("RA-L\u2192RA-R", {"size": 10.5, "bold": True, "color": RED}),
           (" (laterality) \u2014 all scored as full misses.",
            {"size": 10.5, "color": RED})]],
         l=0.45, t=5.85, w=5.75, h=0.5, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)

    # ── RIGHT: the idea ──────────────────────────────────────────────────────
    text(slide, [[("Model the 24 angio_run labels as an aorta-rooted tree "
                   "(branch-of, laterality-mirror, co-opacified-with edges); "
                   "anchor to FMA / RadLex.",
                   {"size": 11, "color": GREY})]],
         l=6.6, t=1.56, w=6.45, h=0.5, anchor=MSO_ANCHOR.TOP)

    cw, ch, x0, gap = 6.45, 0.86, 6.6, 0.99
    card(slide, x0, 2.15, cw, ch, "A", BLUE,
         "Ontology as the metric.",
         "Hierarchical (Wu\u2013Palmer / info-content) + set-valued "
         "(tree-Jaccard) scoring \u2192 partial credit for near hits.")
    card(slide, x0, 2.15 + gap, cw, ch, "B", GREEN,
         "Graph-smoothed kNN + re-rank.",
         "Propagate votes along the tree; constrain by report / series text. "
         "No retraining \u2014 the fastest win.")
    card(slide, x0, 2.15 + 2 * gap, cw, ch, "C", AMBER,
         "Hierarchy-aware embeddings.",
         "Taxonomic-margin or hyperbolic head on frozen RAD-DINO; "
         "tree-constrained multi-label posteriors as the retrieval vector.")
    card(slide, x0, 2.15 + 3 * gap, cw, ch, "D", PURPLE,
         "Temporal + structural.",
         "Model proximal\u2192distal opacification order; scene-graph / GNN "
         "retrieval for multi-vessel runs (novel tier).")

    # ── Bottom strip: novelty + prior work ───────────────────────────────────
    box(slide, 0, 6.45, 13.333, 1.05, fill=LIGHT)
    box(slide, 0, 6.43, 13.333, 0.035, fill=BLUE)
    text(slide,
         [[("Why it\u2019s novel:  ", {"size": 12, "bold": True, "color": RED}),
           ("ontology-guided CBIR for visceral DSA with set-valued, "
            "temporally-aware retrieval \u2014 this combination appears "
            "unexplored.", {"size": 12, "color": DARK})]],
         l=0.45, t=6.55, w=12.5, h=0.4, anchor=MSO_ANCHOR.MIDDLE)
    text(slide,
         [[("Builds on:  ", {"size": 10.5, "bold": True, "color": DARK}),
           ("Bertinetto+ 2020 (better mistakes) \u00b7 Barz & Denzler 2019 "
            "(hierarchy embeddings) \u00b7 Poincar\u00e9 / hyperbolic embeddings "
            "2017\u20132020 \u00b7 C-HMCNN 2020 \u00b7 IRMA (ontology-coded "
            "medical retrieval).", {"size": 10.5, "color": GREY})]],
         l=0.45, t=6.97, w=12.5, h=0.4, anchor=MSO_ANCHOR.MIDDLE)

    out = Path(__file__).resolve().parent / "ontology_retrieval_slide.pptx"
    prs.save(str(out))
    print(f"Saved -> {out}")


if __name__ == "__main__":
    build()
